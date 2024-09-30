from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain.schema.document import Document
from vectordb_handler import load_vectordb
from utils import load_config, timeit
import pypdfium2
import docx
import openpyxl
from pptx import Presentation
import os

config = load_config()
def get_pdf_texts(folder):
    pdf_files = [os.path.join(folder, f) for f in os.listdir(folder) if f.endswith('.pdf')]
    return [extract_text_from_pdf(pdf_file) for pdf_file in pdf_files]

def extract_text_from_pdf(pdf_file):
    with open(pdf_file, "rb") as f:
        pdf_file = pypdfium2.PdfDocument(f)
        return "\n".join(pdf_file.get_page(page_number).get_textpage().get_text_range() for page_number in range(len(pdf_file)))

def get_docx_texts(folder):
    docx_files = [os.path.join(folder, f) for f in os.listdir(folder) if f.endswith('.docx')]
    return [extract_text_from_docx(docx_file) for docx_file in docx_files]

def extract_text_from_docx(docx_file):
    doc = docx.Document(docx_file)
    return "\n".join(paragraph.text for paragraph in doc.paragraphs)

def get_xlsx_texts(folder):
    xlsx_files = [os.path.join(folder, f) for f in os.listdir(folder) if f.endswith('.xlsx')]
    return [extract_text_from_xlsx(xlsx_file) for xlsx_file in xlsx_files]

def extract_text_from_xlsx(xlsx_file):
    workbook = openpyxl.load_workbook(xlsx_file)
    text = []
    for sheet in workbook:
        for row in sheet.iter_rows(values_only=True):
            row_text = " ".join([str(cell) if cell is not None else "" for cell in row])
            text.append(row_text)
    return "\n".join(text)

def get_pptx_texts(folder):
    pptx_files = [os.path.join(folder, f) for f in os.listdir(folder) if f.endswith('.pptx')]
    return [extract_text_from_pptx(pptx_file) for pptx_file in pptx_files]

def extract_text_from_pptx(pptx_file):
    prs = Presentation(pptx_file)
    text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)
    return "\n".join(text)

def get_text_chunks(text):
    splitter = RecursiveCharacterTextSplitter(chunk_size=config["pdf_text_splitter"]["chunk_size"], 
                                              chunk_overlap=config["pdf_text_splitter"]["overlap"],
                                              separators=config["pdf_text_splitter"]["separators"])
    return splitter.split_text(text)

def get_document_chunks(text_list):
    documents = []
    for text in text_list:
        for chunk in get_text_chunks(text):
            documents.append(Document(page_content=chunk))
    return documents

@timeit
def add_documents_to_db():

    pdf_texts = get_pdf_texts(config["folders"]["pdfs"])
    docx_texts = get_docx_texts(config["folders"]["docxs"])
    xlsx_texts = get_xlsx_texts(config["folders"]["xlsxs"])
    pptx_texts = get_pptx_texts(config["folders"]["pptxs"])

    all_texts = pdf_texts + docx_texts + xlsx_texts + pptx_texts

    documents = get_document_chunks(all_texts)
    vector_db = load_vectordb()
    vector_db.add_documents(documents)

    print("Documents from all folders added to db.")
